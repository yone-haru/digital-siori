# -*- coding: utf-8 -*-
"""Yondle A4 poster -> PowerPoint (single A4 portrait slide)."""
import os
from pathlib import Path
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
LP = os.path.join(ROOT, "web", "public", "lp")
SCRATCH = os.path.dirname(os.path.abspath(__file__))

INK = RGBColor(0x0A, 0x0A, 0x0A)
MUTED = RGBColor(0x6E, 0x6B, 0x65)
MUTED2 = RGBColor(0x9A, 0x96, 0x8F)
LINE = RGBColor(0xE3, 0xDF, 0xD6)
PAPER = RGBColor(0xF7, 0xF5, 0xEF)

SERIF_JA = "Yu Mincho Demibold"   # 和文見出し
SERIF_JA_RPR = "游明朝 Demibold"
SERIF_EN = "Georgia"              # 欧文セリフ（Cormorant代替）
GOTHIC = "Yu Gothic Medium"       # 本文
GOTHIC_RPR = "游ゴシック Medium"

DPI = 300
def mm2px(mm): return int(round(mm / 25.4 * DPI))

# ---------- 1) Pillow preprocessing ----------
def rounded(img, radius_px):
    mask = Image.new("L", img.size, 0)
    d = ImageDraw.Draw(mask)
    d.rounded_rectangle([0, 0, img.size[0] - 1, img.size[1] - 1], radius=radius_px, fill=255)
    out = img.convert("RGBA")
    out.putalpha(mask)
    return out

def make_phone(src, dst, inner_w_mm=37.0, pad_mm=1.6, border_mm=0.35):
    img = Image.open(src).convert("RGB")
    iw = mm2px(inner_w_mm)
    ih = int(round(iw * img.size[1] / img.size[0]))
    img = img.resize((iw, ih), Image.LANCZOS)
    img = rounded(img, mm2px(3.6))
    pad, bor = mm2px(pad_mm), mm2px(border_mm)
    W, H = iw + 2 * (pad + bor), ih + 2 * (pad + bor)
    canvas = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(canvas)
    d.rounded_rectangle([0, 0, W - 1, H - 1], radius=mm2px(5), fill=(255, 255, 255, 255),
                        outline=(0xE3, 0xDF, 0xD6, 255), width=bor)
    canvas.paste(img, (pad + bor, pad + bor), img)
    canvas.save(dst)
    return W / DPI * 25.4, H / DPI * 25.4  # size in mm

shots = {}
for name in ["shelf", "timer", "stats"]:
    shots[name] = make_phone(os.path.join(LP, f"{name}.jpg"), os.path.join(SCRATCH, f"phone_{name}.png"))

icon = Image.open(os.path.join(LP, "icon.png")).convert("RGBA").resize((mm2px(11), mm2px(11)), Image.LANCZOS)
rounded(icon, mm2px(2.5)).save(os.path.join(SCRATCH, "icon_rounded.png"))

# ---------- 2) pptx helpers ----------
prs = Presentation()
prs.slide_width = Mm(210)
prs.slide_height = Mm(297)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = PAPER

def set_fonts(run, latin, ea=None, spc_pt=None):
    rPr = run._r.get_or_add_rPr()
    lat = rPr.find(qn("a:latin"))
    if lat is None:
        lat = rPr.makeelement(qn("a:latin"), {})
        rPr.append(lat)
    lat.set("typeface", latin)
    eaEl = rPr.find(qn("a:ea"))
    if eaEl is None:
        eaEl = rPr.makeelement(qn("a:ea"), {})
        rPr.append(eaEl)
    eaEl.set("typeface", ea or latin)
    if spc_pt:
        rPr.set("spc", str(int(spc_pt * 100)))

def textbox(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text,size,color,latin,ea,bold,spc,line_spacing,space_after}"""
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if ln.get("line_spacing"): p.line_spacing = ln["line_spacing"]
        if ln.get("space_after") is not None: p.space_after = Pt(ln["space_after"])
        run = p.add_run()
        run.text = ln["text"]
        f = run.font
        f.size = Pt(ln["size"])
        f.color.rgb = ln.get("color", INK)
        f.bold = ln.get("bold", False)
        set_fonts(run, ln.get("latin", GOTHIC), ln.get("ea"), ln.get("spc"))
    return tb

def hairline(x, y, w, thickness=0.3):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(x), Mm(y), Mm(w), Mm(thickness))
    sh.fill.solid()
    sh.fill.fore_color.rgb = LINE
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

# ---------- 3) layout ----------
MX = 15  # side margin

# Header
slide.shapes.add_picture(os.path.join(SCRATCH, "icon_rounded.png"), Mm(MX), Mm(10), Mm(11), Mm(11))
textbox(29, 9.2, 60, 13, [dict(text="Yondle", size=24, latin=SERIF_EN, spc=1.9)],
        anchor=MSO_ANCHOR.MIDDLE)
textbox(95, 9.2, 100, 13, [dict(text="READING LOG, LIKE A BOOKMARK", size=9, color=MUTED,
        latin=SERIF_EN, spc=2.7)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# Hero
textbox(MX, 27, 170, 6, [dict(text="DIGITAL BOOKSHELF FOR PAPER BOOKS", size=9, color=MUTED,
        latin=SERIF_EN, spc=2.7)])
textbox(MX, 33.5, 175, 32, [
    dict(text="読書の記録を、", size=29, latin=SERIF_JA, ea=SERIF_JA_RPR, spc=1.7, line_spacing=1.35),
    dict(text="しおりのように。", size=29, latin=SERIF_JA, ea=SERIF_JA_RPR, spc=1.7, line_spacing=1.35),
])
textbox(MX, 65.5, 165, 14, [
    dict(text="紙の本の「どこまで読んだか忘れた」をなくす、静かな読書記録アプリ。",
         size=9.5, color=MUTED, ea=GOTHIC_RPR, line_spacing=1.7),
    dict(text="書籍情報の自動取得から、進捗の記録、読書時間の計測、読了予測まで。あなたの読書を、次のページへ。",
         size=9.5, color=MUTED, ea=GOTHIC_RPR, line_spacing=1.7),
])

# Screens
caps = {"shelf": ("LIBRARY", "書影が並ぶ本棚"),
        "timer": ("NOW READING", "没入の読書タイマー"),
        "stats": ("STATS", "積み重ねが見える読書手帖")}
gap = 7
total_w = sum(shots[n][0] for n in shots) + 2 * gap
x = (210 - total_w) / 2
top = 79
for name in ["shelf", "timer", "stats"]:
    w, h = shots[name]
    slide.shapes.add_picture(os.path.join(SCRATCH, f"phone_{name}.png"), Mm(x), Mm(top), Mm(w), Mm(h))
    en, ja = caps[name]
    textbox(x - 5, top + h + 2.5, w + 10, 4, [dict(text=en, size=7.5, color=MUTED2, latin=SERIF_EN, spc=1.9)],
            align=PP_ALIGN.CENTER)
    textbox(x - 5, top + h + 6.5, w + 10, 4, [dict(text=ja, size=7.5, color=MUTED, ea=GOTHIC_RPR)],
            align=PP_ALIGN.CENTER)
    x += w + gap

# Features
feats = [
    ("01", "書影が並ぶ、デジタルの本棚", "タイトル検索だけで表紙・著者・ページ数を自動取得。手元の本が数秒で本棚に並びます。"),
    ("02", "どこまで読んだか、忘れない", "現在のページを記録すれば、進捗率と残りページを常に表示。読書の現在地を見失いません。"),
    ("03", "没入のための、読書タイマー", "黒い画面が集中を守り、読んだ時間とページを静かに記録します。"),
    ("04", "「あと何分で読み終わるか」", "あなた自身の読書ペースから、読了までの残り時間を予測します。"),
    ("05", "積み重ねが見える、読書手帖", "累計読書時間・読了冊数・読んだページ数・連続読書日数を記録します。"),
    ("06", "メモ・評価・タグで整理", "心に残った一文のメモ、評価とレビュー、タグ整理。再読の記録にも対応。"),
]
col_w = 85
fy0, pitch = 183, 23.5
for i, (no, title, body) in enumerate(feats):
    cx = MX if i % 2 == 0 else MX + col_w + 10
    cy = fy0 + (i // 2) * pitch
    hairline(cx, cy, col_w)
    tb = textbox(cx, cy + 2.2, col_w, 7, [dict(text="", size=10.5)])
    p = tb.text_frame.paragraphs[0]
    r1 = p.add_run(); r1.text = no
    r1.font.size = Pt(13); r1.font.color.rgb = MUTED2
    set_fonts(r1, SERIF_EN)
    r2 = p.add_run(); r2.text = "  " + title
    r2.font.size = Pt(10.5); r2.font.color.rgb = INK
    set_fonts(r2, SERIF_JA, SERIF_JA_RPR)
    textbox(cx, cy + 8.6, col_w, 12, [dict(text=body, size=8, color=MUTED, ea=GOTHIC_RPR, line_spacing=1.55)])

# Bottom
by = 262
hairline(MX, by - 4.5, 180, 0.4)
slide.shapes.add_picture(os.path.join(SCRATCH, "..", "..", "..").replace("\\", "/") if False else os.path.join(ROOT, "poster", "qr.png"),
                         Mm(MX), Mm(by), Mm(23), Mm(23))
textbox(46, by + 0.5, 70, 8, [dict(text="無料ではじめる", size=15, latin=SERIF_JA, ea=SERIF_JA_RPR, spc=0.9)])
textbox(46, by + 9, 70, 5, [dict(text="メールアドレスだけで登録できます", size=8, color=MUTED, ea=GOTHIC_RPR)])
textbox(46, by + 15, 70, 7, [dict(text="digital-siori.vercel.app", size=12, latin=SERIF_EN, spc=0.6)])
textbox(115, by + 0.5, 80, 12, [
    dict(text="Web でいますぐ使えます", size=8, color=MUTED, ea=GOTHIC_RPR, line_spacing=1.6),
    dict(text="iOS / Android アプリは近日公開", size=8, color=MUTED, ea=GOTHIC_RPR, line_spacing=1.6),
], align=PP_ALIGN.RIGHT)
textbox(115, by + 16, 80, 5, [dict(text="© 2026 Yondle", size=8, color=MUTED2, latin=SERIF_EN, spc=1.2)],
        align=PP_ALIGN.RIGHT)

out = os.path.join(SCRATCH, "yondle-poster-a4.pptx")
prs.save(out)
print("saved", out)
